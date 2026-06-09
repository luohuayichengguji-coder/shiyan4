from __future__ import annotations

import csv
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.shared import Cm, Inches, Pt, RGBColor
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.enum.section import WD_SECTION
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "outputs"
TABLES = OUT / "tables"
NOTES = OUT / "notes"
FIG = OUT / "figures"

REPORT_MD = OUT / "AI开发技术-实验4.1-医学图像技术调研报告.md"
REPORT_DOCX = OUT / "AI开发技术-实验4.1-医学图像技术调研报告.docx"
PPTX = OUT / "AI开发技术-实验4.1-医学图像技术调研汇报.pptx"


def read_csv(name: str) -> list[dict]:
    with (TABLES / name).open("r", encoding="utf-8-sig") as f:
        return list(csv.DictReader(f))


def read(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def md_table(rows: list[dict], cols: list[str] | None = None, max_rows: int | None = None) -> str:
    if max_rows is not None:
        rows = rows[:max_rows]
    cols = cols or list(rows[0].keys())
    lines = ["| " + " | ".join(cols) + " |", "|" + "|".join(["---"] * len(cols)) + "|"]
    for r in rows:
        lines.append("| " + " | ".join(str(r.get(c, "")).replace("\n", "<br>").replace("|", "\\|") for c in cols) + " |")
    return "\n".join(lines)


def build_markdown() -> str:
    local = read_csv("local_paper_index.csv")
    key = read_csv("key_paper_analysis.csv")
    search = read_csv("literature_search_record.csv")
    sota = read_csv("sota_method_comparison.csv")
    datasets = read_csv("dataset_comparison.csv")
    metrics = read_csv("metrics_comparison.csv")

    refs = [
        "Kirillov A, Mintun E, Ravi N, et al. Segment Anything. arXiv:2304.02643, 2023.",
        "Ma J, He Y, Li F, et al. Segment Anything in Medical Images. arXiv:2304.12306, 2023.",
        "Wu J, Ji W, Liu Y, et al. Medical SAM Adapter: Adapting Segment Anything Model for Medical Image Segmentation. arXiv:2304.12620, 2023.",
        "Cheng J, Ye J, Deng Z, et al. SAM-Med2D. arXiv:2308.16184, 2023.",
        "Ravi N, Gabeur V, Hu Y T, et al. SAM 2: Segment Anything in Images and Videos. arXiv:2408.00714, 2024.",
        "Li W, Xiong X, Xia P, Ju L, Ge Z. TP-DRSeg: Improving Diabetic Retinopathy Lesion Segmentation with Explicit Text-Prompts Assisted SAM. arXiv:2406.15764, 2024.",
        "Huang S, Li J, Xiao Y, Shen N, Xu T. RTNet: Relation Transformer Network for Diabetic Retinopathy Multi-lesion Segmentation. arXiv:2201.11037, 2022.",
        "Jiang H, Gao M, Liu Z, et al. GlanceSeg: Real-time microaneurysm lesion segmentation with gaze-map-guided foundation model for early detection of diabetic retinopathy. arXiv:2311.08075, 2023.",
        "Xiao Q, et al. Improving Lesion Segmentation for Diabetic Retinopathy using Adversarial Learning. arXiv:2007.13854, 2020.",
        "Xia X, Zhan K, Fang Y, Jiang W, Shen F. Lesion-aware network for diabetic retinopathy diagnosis. International Journal of Imaging Systems and Technology, DOI:10.1002/ima.22933.",
        "Wen Y, Luo B, Shi W, et al. SAT-Net: Structure-Aware Transformer-Based Attention Fusion Network for Low-Quality Retinal Fundus Images Enhancement. IEEE Transactions on Multimedia, DOI:10.1109/TMM.2025.3565935.",
        "Patil K D, Palani G, Krishnamurthi G. Efficient Knowledge Distillation of SAM for Medical Image Segmentation. arXiv:2501.16740, 2025.",
        "Zhou C, Li X, Loy C C, Dai B. EdgeSAM: Prompt-In-the-Loop Distillation for SAM. arXiv:2312.06660, 2023.",
        "Moglia A, Leccardi M, Cavicchioli M, et al. Generalist Models in Medical Image Segmentation: A Survey and Performance Comparison with Task-Specific Approaches. arXiv:2506.10825, 2025.",
        "Dai L, Wu L, Li H, et al. A deep learning system for detecting diabetic retinopathy across the disease spectrum. Nature Communications, DOI:10.1038/s41467-021-23458-5.",
        "Dai L, Sheng B, Chen T, et al. A deep learning system for predicting time to progression of diabetic retinopathy. Nature Medicine, DOI:10.1038/s41591-023-02702-z.",
        "Li J, Guan Z, Wang J, et al. Integrated image-based deep learning and language models for primary diabetes care. Nature Medicine, DOI:10.1038/s41591-024-03139-8.",
    ]

    local_counts = {}
    for r in local:
        local_counts[r["category"]] = local_counts.get(r["category"], 0) + 1
    count_text = "；".join(f"{k} {v} 篇" for k, v in local_counts.items())

    lines = [
        "# AI开发技术-实验4.1-医学图像技术调研报告",
        "",
        "## 题目",
        "",
        "基于 SAM 及其医学图像适配方法的眼底视网膜病变/糖尿病视网膜病变病灶分割技术调研",
        "",
        "## 摘要",
        "",
        "本报告依据当前目录下 `AI技术-实验4.1-实验要求-医学图像技术调研.txt` 开展。该任务指导文件已在阶段1首先读取，并作为实验目的、实验内容、实验步骤、实验要求和交付物设计的最高优先级本地依据。围绕眼底视网膜病变/糖尿病视网膜病变病灶分割，本报告结合本地 `眼底视网膜病变分割-论文PDF` 文件夹中的 29 篇论文和 2024-2026 年外部检索文献，系统梳理 SAM、医学 SAM 适配、眼底 DR 病灶分割、低质量眼底增强与 SAM 蒸馏方向。调研发现，通用 SAM 迁移到眼底病灶分割时存在医学领域先验不足、低质量图像影响小病灶和边界、标注稀缺与推理成本较高等问题。报告提出 LQ-Fundus-SAM 方案，即低质量感知增强、眼底专用 SAM 适配、医学先验 prompt、小病灶/边界约束和轻量化蒸馏的组合框架，并给出数据集、训练设置、对比方法、评价指标、消融实验和风险备选方案。报告不编造训练结果，所有性能提升均作为待验证假设。",
        "",
        "关键词：Segment Anything Model；MedSAM；糖尿病视网膜病变；眼底图像；病灶分割；低质量图像增强；知识蒸馏",
        "",
        "## 1 任务依据与 AI 环境配置",
        "",
        "本实验严格对齐任务指导文件。指导文件要求了解医学图像处理、视网膜病变分割和 SAM 技术，使用 Codex 桌面端及相关 Agent/Skill 完成调研、问题分析、解决方案、实验规划、技术报告和 PPT。工作目录为 `/mnt/c/Users/VictorTau/Desktop/hj-5`，所有结果保存在 `outputs/` 下。",
        "",
        "运行环境优先使用 conda 环境 `hj`，命令格式为 `/home/tau/anaconda3/bin/conda run -n hj <command>`。实际检查结果为 Python 3.12.13，PyMuPDF、python-docx、python-pptx、pandas、matplotlib 均可用。过程记录见 `outputs/AI_environment_and_skill_trace.md`。",
        "",
        "本次使用或参考的 skill 包括：academic-research-suite、nature-academic-search、nature-reader、technical-report-writer、nature-writing、nature-figure、nature-paper2ppt、presentations 和 documents。nature-academic-search 的批量搜索接口在本会话中出现事件循环错误，因此采用网页线索加单篇 arXiv/DOI 核验的降级流程，并在检索记录中如实说明。",
        "",
        "## 2 研究背景与意义",
        "",
        "糖尿病视网膜病变是糖尿病患者常见微血管并发症，也是可预防失明的重要原因之一。临床眼底筛查需要识别微动脉瘤、出血、硬性渗出、软性渗出等病灶。图像级 DR 分级可以提示疾病风险，但像素级病灶分割能够解释病变位置、辅助医生复核、支持病灶面积和进展量化，因此具有更强的可解释性和科研价值。",
        "",
        "SAM 的提出使“给定提示即可分割任意对象”成为通用视觉基础模型的重要方向。MedSAM、Medical SAM Adapter、SAM-Med2D 和 Medical SAM 2 等工作进一步说明，SAM 迁移到医学图像需要医学数据、参数高效适配和 prompt 策略。眼底 DR 病灶分割恰好处于通用基础模型和专科医学应用的交汇处：一方面需要 SAM 的泛化和交互能力，另一方面又必须处理眼底图像低质量、小病灶、类别不平衡和边界模糊等细粒度问题。",
        "",
        "![技术路线图](figures/technical_roadmap.png)",
        "",
        "## 3 本地论文库整理",
        "",
        f"本地 PDF 库共抽取 29 篇论文，类别分布为：{count_text}。索引脚本使用 PyMuPDF 自动读取题名、页数、arXiv/DOI、摘要/引言片段、数据集和指标线索，并按 SAM 相关、医学 SAM、眼底病变分割、医学图像增强、SAM 蒸馏等方向建立索引。完整索引见 `outputs/tables/local_paper_index.csv` 与 `outputs/tables/local_paper_index.md`。",
        "",
        "### 3.1 本地论文方向索引摘要",
        "",
        md_table(local, ["id", "year", "category", "direction", "title", "arxiv", "doi"], max_rows=14),
        "",
        "### 3.2 重点论文人工分析",
        "",
        md_table(key, ["ID", "论文", "研究问题", "方法", "数据集/验证", "指标", "优势", "不足", "可借鉴点"]),
        "",
        "## 4 在线检索与最新进展",
        "",
        "阶段4围绕 2024-2026 年 SAM 医学图像分割、眼底病灶分割、低质量眼底增强和轻量化/蒸馏 SAM 检索。关键外部结果包括 TP-DRSeg、Medical SAM 2、MedSAM2、KD-SAM、SAT-Net、AMIR、DiffCode 和 2025 年医学分割 generalist model 综述。完整检索式、来源、日期和筛选理由见 `outputs/notes/literature_search_record.md`。",
        "",
        md_table(search, ["检索日期", "检索式", "来源", "文献", "年份", "标识/链接", "筛选理由"]),
        "",
        "## 5 国内外研究现状与 SOTA 方法分析",
        "",
        "现有研究可归纳为五个方法族：通用 SAM、医学 SAM 适配、眼底 DR 病灶分割、低质量医学/眼底增强和 SAM 蒸馏/轻量化。通用 SAM 提供基础分割范式；医学 SAM 适配解决自然-医学域差距；眼底专用模型强调病灶-血管关系和小病灶检测；低质量增强改善可见性；蒸馏方向降低部署成本。",
        "",
        "![SOTA方法族地图](figures/sota_method_map.png)",
        "",
        md_table(sota),
        "",
        "## 6 常用数据集与评价指标",
        "",
        "数据集需要区分像素级病灶分割数据和图像级 DR 分级数据。IDRiD、DDR、FGADR 更适合病灶分割；APTOS、EyePACS、Messidor 更适合图像级分类、预训练、质量评估或外部筛查泛化讨论。评价指标不能只看 accuracy，应同时报告 per-class Dice/IoU、AUPR、Sensitivity、Specificity、F1、Boundary F1，以及增强和效率指标。",
        "",
        "![数据集与指标选择逻辑](figures/dataset_metric_matrix.png)",
        "",
        "### 6.1 数据集对比",
        "",
        md_table(datasets),
        "",
        "### 6.2 指标对比",
        "",
        md_table(metrics),
        "",
        "## 7 研究不足分析",
        "",
        "![研究不足与方案映射](figures/problem_solution_map.png)",
        "",
        "第一，SAM 从通用图像迁移到眼底 DR 病灶分割时领域适配不足。SAM 训练源主要是自然图像，缺少眼底血管、视盘、病灶形态和 DR 类别先验。MedSAM、Medical SAM Adapter 和 SAM-Med2D 证明医学微调有效，但通用医学 ROI 与 DR 小病灶之间仍有差距。TP-DRSeg 的出现进一步说明，DR 病灶分割需要显式医学概念和文本提示。",
        "",
        "第二，低质量眼底图像会导致微小病灶、边界和细长结构分割不稳定。真实筛查图像常见模糊、低对比、过曝、低分辨率和伪影。SAT-Net 表明结构感知增强可改善眼底图像质量并保留血管细节，但增强指标本身不能保证病灶分割提升，因此增强模块应与分割目标联合约束。",
        "",
        "第三，标注稀缺、类别不平衡、小病灶漏检和推理成本高共同限制实际应用。像素级 DR 数据集规模较小，背景像素占比极高，MA/SE 等小目标容易被整体 Dice 掩盖。SAM/MedSAM 推理成本高，若需要人工点/框 prompt，也会增加临床工作流负担。",
        "",
        "## 8 研究方案设计",
        "",
        "本报告提出 LQ-Fundus-SAM：低质量感知增强 + 眼底专用 SAM 适配 + 小病灶/边界约束 + 轻量化蒸馏的组合方案。目标是在 IDRiD、DDR 等数据上提升 DR 多病灶分割效果，特别关注 MA/HE 小病灶召回、边界质量、低质量图像鲁棒性和推理效率。",
        "",
        "方案包含五个模块：质量感知增强模块 QEM、眼底专用 SAM 适配模块 FSA、医学先验 prompt 模块 MPP、小病灶/边界约束模块 SBC、轻量化蒸馏模块 LSD。Baseline 建议包括 U-Net、DeepLabv3+、HEDNet+cGAN、RTNet、SAM zero-shot、MedSAM、SAM-Med2D、Medical SAM Adapter，以及 GlanceSeg/TP-DRSeg 思路。",
        "",
        "总损失函数可设计为：`L = L_seg + λ1 L_boundary + λ2 L_small + λ3 L_quality + λ4 L_distill`。其中 `L_seg` 为 Dice + Focal/Tversky，`L_boundary` 为边界或 Hausdorff 约束，`L_small` 为小病灶重加权或 lesion-level recall，`L_quality` 保持增强图的结构与病灶纹理，`L_distill` 用于大模型到轻量学生的 embedding/mask/边界蒸馏。",
        "",
        "## 9 实验规划",
        "",
        "![实验流程图](figures/experiment_flow.png)",
        "",
        "实验以 IDRiD 为主数据集，DDR 和 FGADR 用于外部验证或联合训练，APTOS/EyePACS/Messidor 用于预训练、质量评估或筛查泛化讨论。预处理包括眼底圆形视野裁剪、黑边去除、分辨率归一化、颜色标准化、低质量模拟和小病灶增强。训练设置建议冻结 SAM 主干，训练 adapter、prompt 生成器、mask head 和增强模块；若显存有限，优先使用 512×512 输入、梯度累积和轻量学生模型。",
        "",
        "消融实验包括去掉 QEM、去掉 adapter、去掉文本/类别 prompt、去掉结构先验、去掉小病灶损失和去掉蒸馏。可视化应展示原图、增强图、ground truth、baseline 和拟提出模型预测，并单独分析低质量样本、小病灶样本和边界复杂样本。完整实验规划见 `outputs/notes/experiment_plan.md`。",
        "",
        "## 10 AI 过程记录与 Agent/Skill 调用",
        "",
        "本次实验按照 academic-research-suite 的阶段化研究思路完成任务拆解；使用 nature-reader 的源文件定位与结构化阅读思想处理本地重点论文；使用 nature-academic-search 对 arXiv/DOI 进行单篇核验；使用 technical-report-writer 建立合规矩阵和证据映射；使用 nature-figure 生成技术路线图、问题-方案图、实验流程图等；使用 documents 和 presentations 生成 Word 报告与 PPTX。详细轨迹见 `outputs/AI_environment_and_skill_trace.md`。",
        "",
        "## 11 心得体会",
        "",
        "本次实验最大的收获是把“读论文”从线性摘要变成了可追溯的工程化过程：先依据任务指导文件明确交付物，再把本地 PDF、在线检索、SOTA 方法、数据集、指标、问题和方案逐层组织。SAM 类基础模型很强，但医学场景不能简单套用；眼底 DR 病灶分割尤其需要尊重小病灶、低质量成像和临床可解释性。另一个体会是，AI 工具的价值不只是生成文字，而是帮助建立索引、合规矩阵、证据链、方案蓝图和可提交文档。对不确定信息必须标注，不能为了让报告“好看”而编造实验结果。",
        "",
        "## 12 结论",
        "",
        "围绕 SAM 医学图像适配和眼底 DR 病灶分割，本报告完成了本地论文库整理、重点论文阅读、在线文献核验、SOTA/数据集/指标总结、问题分析、研究方案和实验规划。结论是：未来可行方向不是单独改一个模型，而是将低质量增强、眼底领域适配、医学概念 prompt、小病灶/边界损失和轻量化蒸馏结合，形成面向真实筛查场景的端到端方案。",
        "",
        "## 参考文献",
        "",
    ]
    lines.extend(f"{i+1}. {ref}" for i, ref in enumerate(refs))
    lines.append("")
    return "\n".join(lines)


def set_cell_shading(cell, fill: str):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_cell_text(cell, text: str, bold=False, size=8):
    cell.text = ""
    p = cell.paragraphs[0]
    r = p.add_run(str(text))
    r.font.name = "Microsoft YaHei"
    r._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    r.font.size = Pt(size)
    r.bold = bold
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def add_table_docx(doc: Document, rows: list[dict], cols: list[str], max_rows: int | None = None):
    if max_rows is not None:
        rows = rows[:max_rows]
    table = doc.add_table(rows=1, cols=len(cols))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    for i, col in enumerate(cols):
        set_cell_text(hdr[i], col, bold=True, size=8)
        set_cell_shading(hdr[i], "E8EEF5")
    for row in rows:
        cells = table.add_row().cells
        for i, col in enumerate(cols):
            set_cell_text(cells[i], row.get(col, ""), size=7)
    doc.add_paragraph()


def setup_doc_styles(doc: Document):
    section = doc.sections[0]
    section.top_margin = Inches(0.8)
    section.bottom_margin = Inches(0.8)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)
    styles = doc.styles
    styles["Normal"].font.name = "Microsoft YaHei"
    styles["Normal"]._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    styles["Normal"].font.size = Pt(10.5)
    for name, size, color in [("Heading 1", 16, "1F4D78"), ("Heading 2", 13, "2E74B5"), ("Heading 3", 11, "1F4D78")]:
        style = styles[name]
        style.font.name = "Microsoft YaHei"
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
        style.font.size = Pt(size)
        style.font.color.rgb = RGBColor.from_string(color)
        style.font.bold = True


def add_para(doc, text: str):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(6)
    p.paragraph_format.line_spacing = 1.15
    run = p.add_run(text)
    run.font.name = "Microsoft YaHei"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    run.font.size = Pt(10.5)


def add_picture(doc, filename: str, caption: str):
    path = FIG / filename
    if path.exists():
        doc.add_picture(str(path), width=Inches(6.6))
        last = doc.paragraphs[-1]
        last.alignment = WD_ALIGN_PARAGRAPH.CENTER
        cap = doc.add_paragraph(caption)
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        cap.runs[0].font.size = Pt(8)
        cap.runs[0].font.color.rgb = RGBColor(90, 90, 90)


def build_docx(md_text: str):
    local = read_csv("local_paper_index.csv")
    key = read_csv("key_paper_analysis.csv")
    search = read_csv("literature_search_record.csv")
    sota = read_csv("sota_method_comparison.csv")
    datasets = read_csv("dataset_comparison.csv")
    metrics = read_csv("metrics_comparison.csv")

    doc = Document()
    setup_doc_styles(doc)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run("AI开发技术-实验4.1-医学图像技术调研报告")
    r.font.name = "Microsoft YaHei"
    r._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = RGBColor(31, 77, 120)
    subtitle = doc.add_paragraph("基于 SAM 及其医学图像适配方法的眼底视网膜病变/糖尿病视网膜病变病灶分割技术调研")
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.runs[0].font.size = Pt(12)
    doc.add_paragraph()

    doc.add_heading("摘要", level=1)
    add_para(doc, "本报告依据当前目录下 AI技术-实验4.1-实验要求-医学图像技术调研.txt 开展。该任务指导文件已在阶段1首先读取，并作为实验目的、实验内容、实验步骤、实验要求和交付物设计的最高优先级本地依据。报告结合本地 29 篇 PDF 和 2024-2026 年外部检索文献，系统梳理 SAM、医学 SAM 适配、眼底 DR 病灶分割、低质量眼底增强与 SAM 蒸馏方向，提出 LQ-Fundus-SAM 组合研究方案和实验规划。")
    add_para(doc, "关键词：Segment Anything Model；MedSAM；糖尿病视网膜病变；眼底图像；病灶分割；低质量图像增强；知识蒸馏")

    doc.add_heading("1 任务依据与环境配置", level=1)
    add_para(doc, "本实验首先读取并引用本地任务指导文件，并将其作为最高优先级任务依据。运行环境优先使用 conda hj，命令统一采用 /home/tau/anaconda3/bin/conda run -n hj <command>。检查结果显示 Python 3.12.13，PyMuPDF、python-docx、python-pptx、pandas、matplotlib 均可用。")
    add_para(doc, "使用或参考的 skill 包括 academic-research-suite、nature-academic-search、nature-reader、technical-report-writer、nature-writing、nature-figure、nature-paper2ppt、presentations 和 documents。详细轨迹见 outputs/AI_environment_and_skill_trace.md。")

    doc.add_heading("2 研究背景与意义", level=1)
    add_para(doc, "糖尿病视网膜病变是糖尿病患者常见微血管并发症，也是可预防失明的重要原因之一。像素级病灶分割能够解释病变位置、辅助医生复核并支持进展量化。SAM 的 promptable segmentation 范式为医学图像分割提供了新工具，但眼底 DR 病灶分割仍需要处理低质量、小目标、类别不平衡和边界模糊问题。")
    add_picture(doc, "technical_roadmap.png", "图1 LQ-Fundus-SAM 技术路线图")

    doc.add_heading("3 本地论文库整理与重点论文分析", level=1)
    add_para(doc, "本地 PDF 库共抽取 29 篇论文，覆盖 SAM 相关、医学图像增强、眼底视网膜病变分割和 SAM 模型蒸馏。完整索引见 outputs/tables/local_paper_index.csv。")
    add_table_docx(doc, local, ["id", "year", "category", "direction", "title", "arxiv"], max_rows=12)
    doc.add_heading("重点论文人工分析", level=2)
    add_table_docx(doc, key, ["ID", "论文", "研究问题", "方法", "优势", "不足", "可借鉴点"])

    doc.add_heading("4 在线检索与最新进展", level=1)
    add_para(doc, "阶段4围绕 2024-2026 年 SAM 医学图像分割、眼底病灶分割、低质量眼底增强和轻量化/蒸馏 SAM 检索。nature-academic-search 批量搜索接口出现事件循环错误，因此采用网页检索线索加单篇 arXiv/DOI 核验的降级流程。")
    add_table_docx(doc, search, ["检索日期", "文献", "年份", "标识/链接", "筛选理由"])

    doc.add_heading("5 SOTA 方法、数据集与指标", level=1)
    add_picture(doc, "sota_method_map.png", "图2 SOTA 方法族地图")
    add_table_docx(doc, sota, ["方向", "代表方法", "核心思想", "主要局限"])
    doc.add_heading("数据集与指标", level=2)
    add_picture(doc, "dataset_metric_matrix.png", "图3 数据集与指标选择逻辑")
    add_table_docx(doc, datasets, ["数据集", "任务属性", "规模/标注", "适用性", "注意事项"])
    add_table_docx(doc, metrics, ["指标", "定义/含义", "适用场景", "注意事项"])

    doc.add_heading("6 研究不足", level=1)
    add_picture(doc, "problem_solution_map.png", "图4 研究不足与方案模块映射")
    add_para(doc, "第一，SAM 从通用图像迁移到眼底 DR 病灶分割时领域适配不足。第二，低质量眼底图像会导致微小病灶、边界和细长结构分割不稳定。第三，标注稀缺、类别不平衡、小病灶漏检和推理成本高共同限制实际应用。")

    doc.add_heading("7 研究方案", level=1)
    add_para(doc, "本报告提出 LQ-Fundus-SAM：低质量感知增强 + 眼底专用 SAM 适配 + 小病灶/边界约束 + 轻量化蒸馏的组合方案。模块包括 QEM、FSA、MPP、SBC 和 LSD。总损失函数为 L = L_seg + λ1 L_boundary + λ2 L_small + λ3 L_quality + λ4 L_distill。")
    add_para(doc, "Baseline 建议包括 U-Net、DeepLabv3+、HEDNet+cGAN、RTNet、SAM zero-shot、MedSAM、SAM-Med2D、Medical SAM Adapter，以及 GlanceSeg/TP-DRSeg 思路。")

    doc.add_heading("8 实验规划", level=1)
    add_picture(doc, "experiment_flow.png", "图5 实验规划流程")
    add_para(doc, "实验以 IDRiD 为主数据集，DDR 和 FGADR 用于外部验证或联合训练，APTOS/EyePACS/Messidor 用于预训练、质量评估或筛查泛化讨论。消融实验包括去掉 QEM、adapter、文本/类别 prompt、结构先验、小病灶损失和蒸馏。")

    doc.add_heading("9 AI 过程记录与心得体会", level=1)
    add_para(doc, "本次实验把读论文、检索、证据表、方案设计和交付物生成组织为可追踪流程。AI 工具的价值不只是生成文字，而是帮助建立索引、合规矩阵、证据链和可提交文档。对不确定信息必须标注，不能编造实验结果。")

    doc.add_heading("10 结论", level=1)
    add_para(doc, "围绕 SAM 医学图像适配和眼底 DR 病灶分割，本报告完成了本地论文库整理、重点论文阅读、在线文献核验、SOTA/数据集/指标总结、问题分析、研究方案和实验规划。可行方向是将低质量增强、眼底领域适配、医学概念 prompt、小病灶/边界损失和轻量化蒸馏结合，形成面向真实筛查场景的端到端方案。")

    doc.add_heading("参考文献", level=1)
    refs = md_text.split("## 参考文献\n\n", 1)[1].strip().splitlines()
    for ref in refs:
        add_para(doc, ref)

    doc.save(REPORT_DOCX)


def add_slide_title(slide, title: str, subtitle: str | None = None):
    tx = slide.shapes.add_textbox(PInches(0.45), PInches(0.25), PInches(12.4), PInches(0.6))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = PPt(26)
    p.font.bold = True
    p.font.color.rgb = PRGBColor(31, 41, 55)
    if subtitle:
        sub = slide.shapes.add_textbox(PInches(0.48), PInches(0.88), PInches(12.0), PInches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Microsoft YaHei"
        sp.font.size = PPt(11)
        sp.font.color.rgb = PRGBColor(107, 114, 128)


def add_bullets(slide, bullets, x, y, w, h, size=16):
    box = slide.shapes.add_textbox(PInches(x), PInches(y), PInches(w), PInches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = PPt(size)
        p.font.color.rgb = PRGBColor(31, 41, 55)
        p.space_after = PPt(8)
    return box


def add_card(slide, x, y, w, h, title, body, fill=(239, 246, 255)):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PInches(x), PInches(y), PInches(w), PInches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = PRGBColor(*fill)
    shp.line.color.rgb = PRGBColor(203, 213, 225)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = PInches(0.12)
    tf.margin_right = PInches(0.12)
    tf.margin_top = PInches(0.10)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = PPt(15)
    p.font.bold = True
    p.font.color.rgb = PRGBColor(31, 41, 55)
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = "Microsoft YaHei"
    p2.font.size = PPt(11)
    p2.font.color.rgb = PRGBColor(55, 65, 81)
    return shp


def add_image(slide, name, x, y, w):
    slide.shapes.add_picture(str(FIG / name), PInches(x), PInches(y), width=PInches(w))


def build_pptx():
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "AI开发技术-实验4.1", "医学图像技术调研")
    tx = slide.shapes.add_textbox(PInches(0.8), PInches(2.0), PInches(11.8), PInches(1.2))
    p = tx.text_frame.paragraphs[0]
    tx.text_frame.word_wrap = True
    p.text = "基于 SAM 及其医学图像适配方法的\n眼底视网膜病变/糖尿病视网膜病变病灶分割技术调研"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Microsoft YaHei"; p.font.size = PPt(28); p.font.bold = True; p.font.color.rgb = PRGBColor(31, 41, 55)
    add_bullets(slide, ["任务依据：AI技术-实验4.1-实验要求-医学图像技术调研.txt", "本地 PDF：29 篇；外部检索：2024-2026 SAM/DR/增强/蒸馏方向", "交付物：过程记录、技术报告 Word、汇报 PPT、README、合规矩阵"], 0.9, 4.1, 11.5, 1.5, 15)

    # 2
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "任务要求与执行路径", "阶段1-10对齐课程指导文件")
    add_card(slide, 0.7, 1.4, 3.7, 1.3, "输入", "任务指导文件 + 本地论文PDF + 在线检索", (239, 246, 255))
    add_card(slide, 4.8, 1.4, 3.7, 1.3, "过程", "论文索引、重点阅读、SOTA、问题分析、方案设计", (240, 253, 250))
    add_card(slide, 8.9, 1.4, 3.7, 1.3, "输出", "报告 DOCX、PPTX、README、合规矩阵、全过程记录", (240, 253, 244))
    add_bullets(slide, ["已使用/参考：academic-research-suite、nature-reader、nature-academic-search。", "报告/图表/文档/PPT：technical-report-writer、nature-figure、documents、presentations。", "所有结果保存在当前 hj-5 工作目录的 outputs/ 下。"], 0.9, 3.45, 11.5, 1.35, 15)

    # 3
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "研究背景", "DR病灶分割需要解释性、像素级定位和临床可用性")
    add_card(slide, 0.7, 1.3, 3.8, 1.5, "临床意义", "DR 是糖尿病重要眼部并发症；早筛可降低视力损失风险。", (254, 242, 242))
    add_card(slide, 4.75, 1.3, 3.8, 1.5, "技术价值", "像素级分割可定位 MA、HE、EX、SE，支持解释和量化。", (239, 246, 255))
    add_card(slide, 8.8, 1.3, 3.8, 1.5, "核心矛盾", "SAM 泛化强，但眼底小病灶、低质量和类别不平衡仍难。", (254, 243, 199))
    add_image(slide, "dataset_metric_matrix.png", 1.0, 3.25, 11.2)

    # 4
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "本地论文库整理", "29篇PDF按五个方向建立索引")
    add_image(slide, "sota_method_map.png", 0.8, 1.15, 7.2)
    add_bullets(slide, ["SAM相关：SAM、SAM2、HQ-SAM、FastSAM、Semantic-SAM", "医学SAM：MedSAM、Medical SAM Adapter、SAM-Med2D", "眼底DR：HEDNet+cGAN、RTNet、GlanceSeg、LANet、DeepLabv3+", "增强/蒸馏：SAT-Net、AMIR、DiffCode、KD-SAM、EdgeSAM"], 8.25, 1.35, 4.4, 4.6, 15)

    # 5
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "重点论文证据链", "从通用SAM到眼底DR专用适配")
    add_card(slide, 0.6, 1.25, 3.0, 1.25, "SAM / SAM2", "promptable segmentation；图像/视频统一与记忆机制", (219, 234, 254))
    add_card(slide, 3.85, 1.25, 3.0, 1.25, "MedSAM / SAM-Med2D", "大规模医学微调与多prompt策略，缩小医学域差距", (204, 251, 241))
    add_card(slide, 7.1, 1.25, 3.0, 1.25, "RTNet / GlanceSeg", "病灶-血管关系、gaze/saliency prompt 处理小病灶", (220, 252, 231))
    add_card(slide, 10.35, 1.25, 2.4, 1.25, "SAT-Net / KD-SAM", "低质量增强与医学SAM蒸馏", (237, 233, 254))
    add_bullets(slide, ["结论：单一路线不足，组合式方案更符合任务问题。", "关键缺口：领域适配、低质量、小病灶、效率。"], 1.0, 3.4, 11.2, 1.2, 18)

    # 6
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "最新检索结果", "2024-2026 关键补充论文")
    add_bullets(slide, ["TP-DRSeg（2024）：显式文本提示辅助 SAM 做 DR 病灶分割。", "Medical SAM 2 / MedSAM2：SAM2 向医学图像、3D和视频扩展。", "KD-SAM（2025）：医学 SAM 轻量化与蒸馏。", "SAT-Net（2025）：低质量眼底增强、结构保持和下游任务验证。", "Generalist medical segmentation survey（2025）：总结 SAM、adapter、fine-tuning 和临床转化挑战。"], 0.8, 1.4, 11.8, 4.5, 17)

    # 7
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "研究不足", "三个瓶颈决定方案设计")
    add_image(slide, "problem_solution_map.png", 0.8, 1.15, 11.7)

    # 8
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "拟提出方案：LQ-Fundus-SAM", "低质量感知增强 + 眼底专用SAM适配 + 小病灶/边界约束 + 蒸馏")
    add_image(slide, "technical_roadmap.png", 0.7, 1.05, 12.0)

    # 9
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "模型模块与损失函数", "让增强、分割和蒸馏服务同一个目标")
    add_bullets(slide, ["QEM：质量感知增强，保留血管和病灶纹理。", "FSA：眼底专用 adapter/LoRA 和 lesion-aware token。", "MPP：自动点/框/热图 + 医学文本类别 prompt。", "SBC：Dice + Focal/Tversky + Boundary + 小病灶重加权。", "LSD：teacher-student 蒸馏，约束 embedding、mask、边界和小病灶响应。"], 0.8, 1.25, 6.0, 4.9, 16)
    add_card(slide, 7.2, 2.0, 5.2, 1.4, "总损失", "L = L_seg + λ1 L_boundary + λ2 L_small + λ3 L_quality + λ4 L_distill", (239, 246, 255))
    add_card(slide, 7.2, 3.8, 5.2, 1.4, "不编造结果", "本报告给出可执行蓝图，所有性能提升均为待实验验证假设。", (254, 242, 242))

    # 10
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "实验规划", "数据、训练、对比、消融与风险控制")
    add_image(slide, "experiment_flow.png", 0.7, 1.1, 12.0)

    # 11
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "评价指标与消融", "小病灶不能只看整体Dice")
    add_card(slide, 0.7, 1.35, 3.7, 1.4, "分割效果", "per-class Dice/IoU、AUPR、Sensitivity、Specificity、F1、Boundary F1", (239, 246, 255))
    add_card(slide, 4.85, 1.35, 3.7, 1.4, "增强质量", "PSNR、SSIM、LPIPS，同时报告下游分割收益", (240, 253, 250))
    add_card(slide, 9.0, 1.35, 3.4, 1.4, "部署效率", "FPS、参数量、FLOPs、显存占用", (237, 233, 254))
    add_bullets(slide, ["消融：去QEM、去adapter、去文本prompt、去结构先验、去小病灶损失、去蒸馏。", "可视化：低质量样本、微小病灶样本、边界复杂样本和错误案例。"], 0.9, 3.6, 11.8, 1.4, 17)

    # 12
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "总结与心得", "从论文堆到可执行研究蓝图")
    add_bullets(slide, ["本地任务指导文件已作为全过程最高优先级依据。", "本地29篇PDF + 外部检索共同支撑调研，避免空泛综述。", "SAM 医学适配不能只做模型替换，必须结合领域先验、图像质量、小病灶损失和部署效率。", "AI工具最有价值的部分是建立证据链、合规矩阵和可复查交付物。", "后续工作：按实验规划下载/整理数据，复现baseline，再逐步验证LQ-Fundus-SAM各模块。"], 0.9, 1.35, 11.7, 4.7, 18)

    prs.save(PPTX)


def main() -> None:
    md = build_markdown()
    REPORT_MD.write_text(md, encoding="utf-8")
    build_docx(md)
    build_pptx()
    print(REPORT_MD)
    print(REPORT_DOCX)
    print(PPTX)


if __name__ == "__main__":
    main()
